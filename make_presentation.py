"""Generate a PowerPoint presentation about the Flappy Bird World Model project.
Focus: design choices, failures, iterative debugging journey."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Color palette
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_MED = RGBColor(0x22, 0x22, 0x3A)
ACCENT_BLUE = RGBColor(0x4E, 0xA8, 0xDE)
ACCENT_GREEN = RGBColor(0x4E, 0xDE, 0x8A)
ACCENT_ORANGE = RGBColor(0xDE, 0x8A, 0x4E)
ACCENT_PURPLE = RGBColor(0x9B, 0x6E, 0xDE)
ACCENT_RED = RGBColor(0xDE, 0x5E, 0x5E)
ACCENT_YELLOW = RGBColor(0xDE, 0xDE, 0x4E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_RED_BG = RGBColor(0x3E, 0x1A, 0x1A)
DARK_GREEN_BG = RGBColor(0x1A, 0x3E, 0x1A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_slide_bg(slide, color=BG_DARK):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=LIGHT_GRAY, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def add_code_block(slide, left, top, width, height, code, font_size=11):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    fill = txBox.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0x12, 0x12, 0x20)
    for i, line in enumerate(code.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(0xA0, 0xD0, 0xA0)
        p.font.name = "Courier New"
        p.space_after = Pt(2)
    return txBox


def add_box(slide, left, top, width, height, fill_color=BG_MED, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT_BLUE):
    connector = slide.shapes.add_connector(
        1, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector


def add_section_divider(title, subtitle="", color=ACCENT_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.2), Inches(13.333), Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_textbox(slide, 1.5, 2.0, 10, 1.0, title, font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, 1.5, 3.6, 10, 0.8, subtitle, font_size=20, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
    return slide


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 1.5, 1.5, 10, 1.5, "Flappy Bird World Model", font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.2, 10, 1.0, "Learning Game Physics with Flow Matching", font_size=28, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 4.5, 10, 0.8, "A story about design choices, debugging failures,\nand the surprisingly hard problem of making a bird fall down", font_size=16, color=MED_GRAY, alignment=PP_ALIGN.CENTER)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.2), Inches(5.3), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

# ============================================================
# SLIDE 2: Project Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.4, 12, 0.8, "Project Overview", font_size=36, color=WHITE, bold=True)

add_box(slide, 0.8, 1.5, 5.5, 2.2, border_color=ACCENT_BLUE)
add_textbox(slide, 1.0, 1.6, 5.0, 0.5, "Goal", font_size=22, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, 1.0, 2.2, 5.0, 1.5, [
    "Train a neural network to simulate Flappy Bird entirely",
    "Given context frames + player action, predict next frame",
    "Interactive play — no game engine at inference",
], font_size=14, spacing=Pt(6))

add_box(slide, 7.0, 1.5, 5.5, 2.2, border_color=ACCENT_GREEN)
add_textbox(slide, 7.2, 1.6, 5.0, 0.5, "Approach vs GameNGen", font_size=22, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, 7.2, 2.2, 5.0, 1.5, [
    "Inspired by GameNGen (Google, 2024) — but built from scratch",
    "Flow matching instead of DDIM diffusion",
    "~3.4M params vs GameNGen's ~860M (Stable Diffusion 1.4)",
    "~470K frames vs 900M frames — 2000x less data",
], font_size=14, spacing=Pt(6))

add_box(slide, 0.8, 4.0, 12.0, 3.2, border_color=ACCENT_PURPLE)
add_textbox(slide, 1.0, 4.1, 11.0, 0.5, "Pipeline", font_size=22, color=ACCENT_PURPLE, bold=True)

stages = [
    ("PPO Agent", "game/rl/", ACCENT_BLUE),
    ("VOD Recording\n(hijacked agent)", "vod/", ACCENT_GREEN),
    ("VAE Training", "diffuse/vae/", ACCENT_ORANGE),
    ("Latent Encoding", "latent-vod/", ACCENT_PURPLE),
    ("Flow Model", "diffuse/ngen/", ACCENT_RED),
    ("Interactive\nInference", "world/", ACCENT_YELLOW),
]
for i, (title, path, color) in enumerate(stages):
    x = 1.0 + i * 2.0
    add_box(slide, x, 4.8, 1.7, 1.6, border_color=color)
    add_textbox(slide, x + 0.05, 4.85, 1.6, 0.7, title, font_size=11, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.05, 5.7, 1.6, 0.4, path, font_size=9, color=MED_GRAY, alignment=PP_ALIGN.CENTER, font_name="Courier New")
    if i < len(stages) - 1:
        add_arrow(slide, x + 1.7 + 0.05, 5.6, x + 2.0 - 0.05, 5.6, color=MED_GRAY)

add_textbox(slide, 1.0, 6.6, 11.0, 0.4, "467K frames | 3,135 episodes | 17 hijacking configs | 4.3hr gameplay | 2x H200 GPUs", font_size=13, color=MED_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3: Flow Matching in 60 seconds
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.4, 12, 0.8, "Flow Matching in 60 Seconds", font_size=36, color=WHITE, bold=True)

add_box(slide, 0.8, 1.5, 5.8, 2.5, border_color=ACCENT_BLUE)
add_textbox(slide, 1.0, 1.6, 5.4, 0.5, "Training", font_size=22, color=ACCENT_BLUE, bold=True)
add_code_block(slide, 1.0, 2.2, 5.4, 1.6,
    "z_0 ~ N(0,1)            # random noise\n"
    "z_1 = target_latent     # real game frame\n"
    "t ~ U(0,1)              # random time\n"
    "z_t = (1-t)*z_0 + t*z_1 # interpolate\n"
    "v_target = z_1 - z_0    # target velocity\n"
    "loss = MSE(model(z_t, t, cond), v_target)",
    font_size=12)

add_box(slide, 7.0, 1.5, 5.8, 2.5, border_color=ACCENT_GREEN)
add_textbox(slide, 7.2, 1.6, 5.4, 0.5, "Inference (50 Euler steps)", font_size=22, color=ACCENT_GREEN, bold=True)
add_code_block(slide, 7.2, 2.2, 5.4, 1.6,
    "z = N(0,1)              # start from noise\n"
    "for i in range(50):\n"
    "    t = (i + 0.5) / 50  # midpoint rule\n"
    "    v = model(z, t, cond)\n"
    "    z = z + v * dt\n"
    "    z = clamp(z, -4, 4) # stability",
    font_size=12)

add_box(slide, 0.8, 4.3, 12.0, 2.8, border_color=ACCENT_ORANGE)
add_textbox(slide, 1.0, 4.4, 11.0, 0.5, "Key Idea: Learn a velocity field that transports noise to data along straight paths", font_size=18, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, 1.0, 5.0, 5.5, 1.8, [
    "Simpler than denoising diffusion — no noise schedule needed",
    "Straight-line paths = Euler integrator is a good approximation",
    "Loss is plain MSE on velocity predictions",
    "Connects naturally to reflow for fewer inference steps",
], font_size=14, spacing=Pt(4))
add_bullet_list(slide, 7.0, 5.0, 5.5, 1.8, [
    "Why vs DDIM? Simpler math, no variance weighting",
    "Midpoint rule: fixes systematic bias (ODE never reached t=1.0)",
    "Latent clamping: prevents ODE drift to out-of-distribution",
    "CFG: 2 forward passes per step (100 evals per frame)",
], font_size=14, spacing=Pt(4))


# ============================================================
# SLIDE 4: Architecture at a Glance
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.4, 12, 0.8, "ResUNet Architecture", font_size=36, color=WHITE, bold=True)

# Architecture diagram
add_box(slide, 0.5, 1.5, 2.5, 4.5, border_color=ACCENT_BLUE)
add_textbox(slide, 0.6, 1.6, 2.3, 0.5, "Encoder", font_size=20, color=ACCENT_BLUE, bold=True)
for i, (name, dims, chs) in enumerate([
    ("DownRes 1", "36x64 -> 18x32", "36ch -> 64ch"),
    ("DownRes 2", "18x32 -> 9x16", "64ch -> 128ch"),
]):
    y = 2.3 + i * 1.3
    add_box(slide, 0.7, y, 2.1, 1.0, fill_color=RGBColor(0x1E, 0x3A, 0x5F))
    add_textbox(slide, 0.8, y + 0.05, 1.9, 0.35, name, font_size=13, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.8, y + 0.35, 1.9, 0.3, dims, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 0.8, y + 0.6, 1.9, 0.3, chs, font_size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_box(slide, 3.5, 2.0, 2.8, 3.5, border_color=ACCENT_PURPLE)
add_textbox(slide, 3.6, 2.1, 2.6, 0.5, "Bottleneck", font_size=20, color=ACCENT_PURPLE, bold=True)
add_box(slide, 3.7, 2.7, 2.4, 0.7, fill_color=RGBColor(0x2E, 0x1E, 0x5F))
add_textbox(slide, 3.8, 2.75, 2.2, 0.5, "ResBlock\n128ch -> 256ch", font_size=12, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)
add_box(slide, 3.7, 3.6, 2.4, 0.7, fill_color=RGBColor(0x2E, 0x1E, 0x5F))
add_textbox(slide, 3.8, 3.65, 2.2, 0.5, "Self-Attention\n8 heads, 144 tokens", font_size=12, color=ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)
add_box(slide, 3.7, 4.5, 2.4, 0.7, fill_color=RGBColor(0x3E, 0x1E, 0x1E))
add_textbox(slide, 3.8, 4.55, 2.2, 0.5, "Done head (detach)\nDynamics head", font_size=11, color=ACCENT_RED, alignment=PP_ALIGN.CENTER)

add_box(slide, 6.8, 1.5, 2.5, 4.5, border_color=ACCENT_GREEN)
add_textbox(slide, 6.9, 1.6, 2.3, 0.5, "Decoder", font_size=20, color=ACCENT_GREEN, bold=True)
for i, (name, dims, chs) in enumerate([
    ("UpRes 1", "9x16 -> 18x32", "256ch -> 128ch"),
    ("UpRes 2", "18x32 -> 36x64", "128ch -> 64ch"),
]):
    y = 2.3 + i * 1.3
    add_box(slide, 7.0, y, 2.1, 1.0, fill_color=RGBColor(0x1E, 0x3A, 0x2F))
    add_textbox(slide, 7.1, y + 0.05, 1.9, 0.35, name, font_size=13, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 7.1, y + 0.35, 1.9, 0.3, dims, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 7.1, y + 0.6, 1.9, 0.3, chs, font_size=11, color=MED_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 3.1, 2.7, 0.5, 0.4, "-->", font_size=14, color=MED_GRAY)
add_textbox(slide, 3.1, 4.0, 0.5, 0.4, "-->", font_size=14, color=MED_GRAY)
add_textbox(slide, 6.3, 2.7, 0.5, 0.4, "-->", font_size=14, color=MED_GRAY)
add_textbox(slide, 6.3, 4.0, 0.5, 0.4, "-->", font_size=14, color=MED_GRAY)

# Conditioning panel
add_box(slide, 9.8, 1.5, 3.0, 5.0, border_color=ACCENT_ORANGE)
add_textbox(slide, 9.9, 1.6, 2.8, 0.5, "Conditioning", font_size=20, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, 9.9, 2.2, 2.8, 4.0, [
    "Input: z_t + z_cond (k=4 context) + action spatial (16ch)",
    "",
    "Time: sinusoidal -> SiLU -> Linear -> 128d",
    "",
    "Aug: Embedding(10, 128) -> SiLU",
    "",
    "Actions: all k+1 -> Embed -> flatten -> Linear -> 128d",
    "",
    "AdaGN: cat([t, c, aug]) -> scale/shift at every ResBlock",
    "",
    "Spatial: target action broadcast to (B,16,H,W) at input",
], font_size=10, spacing=Pt(2))


# ============================================================
# SECTION DIVIDER: THE JOURNEY
# ============================================================
add_section_divider("The Journey", "Bugs, failures, and design decisions along the way", ACCENT_RED)


# ============================================================
# SLIDE: Bug #1 — ODE Time Stepping
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.3, 8, 0.6, "Bug: ODE Never Reaches t=1.0", font_size=32, color=ACCENT_RED, bold=True)
add_textbox(slide, 9.5, 0.4, 3, 0.4, "Severity: CRITICAL", font_size=16, color=ACCENT_RED, bold=True)

add_box(slide, 0.8, 1.2, 5.5, 2.8, fill_color=DARK_RED_BG, border_color=ACCENT_RED)
add_textbox(slide, 1.0, 1.3, 5.0, 0.4, "The Bug", font_size=20, color=ACCENT_RED, bold=True)
add_code_block(slide, 1.0, 1.8, 5.0, 1.8,
    "# BUGGY: left-endpoint Euler\n"
    "for i in range(50):  # num_steps=50\n"
    "    t = i * dt       # t = 0, 0.02, ..., 0.98\n"
    "    v = model(z, t)  # NEVER evaluates at t=1.0!\n"
    "    z = z + v * dt\n"
    "\n"
    "# Systematic 2% trajectory error\n"
    "# Forward/backward ODE asymmetry breaks reflow",
    font_size=11)

add_box(slide, 7.0, 1.2, 5.5, 2.8, fill_color=DARK_GREEN_BG, border_color=ACCENT_GREEN)
add_textbox(slide, 7.2, 1.3, 5.0, 0.4, "The Fix: Midpoint Rule", font_size=20, color=ACCENT_GREEN, bold=True)
add_code_block(slide, 7.2, 1.8, 5.0, 1.8,
    "# FIXED: midpoint Euler\n"
    "for i in range(50):\n"
    "    t = (i + 0.5) * dt  # t = 0.01, 0.03, ..., 0.99\n"
    "    v = model(z, t)\n"
    "    z = z + v * dt\n"
    "\n"
    "# Centers evaluation in each interval\n"
    "# Forward/backward become symmetric",
    font_size=11)

add_box(slide, 0.8, 4.3, 12.0, 2.8, border_color=ACCENT_ORANGE)
add_textbox(slide, 1.0, 4.4, 11.0, 0.5, "Why It Matters for Reflow", font_size=20, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, 1.0, 5.0, 5.5, 1.5, [
    "Reflow needs backward ODE: z_0 = backward(z_1)",
    "If forward ends at t=0.98, backward starts at t=1.0",
    "Asymmetric integration = paired points aren't truly paired",
    "Training on bad pairs makes reflow model worse, not better",
], font_size=14, spacing=Pt(4))
add_bullet_list(slide, 7.0, 5.0, 5.5, 1.5, [
    "Also found: train used 50 steps, inference used 20 steps",
    "2.5x larger dt at inference = different discretization error",
    "Fix: always use 50 steps for both",
    "Midpoint rule also reduces error at any step count",
], font_size=14, spacing=Pt(4))


# ============================================================
# SLIDE: Bug #2 — CFG Dropping Actions
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.3, 8, 0.6, "Bug: CFG Drops Actions — Bird Disappears", font_size=32, color=ACCENT_RED, bold=True)
add_textbox(slide, 9.5, 0.4, 3, 0.4, "Severity: CRITICAL", font_size=16, color=ACCENT_RED, bold=True)

add_box(slide, 0.8, 1.2, 5.5, 3.2, fill_color=DARK_RED_BG, border_color=ACCENT_RED)
add_textbox(slide, 1.0, 1.3, 5.0, 0.4, "The Bug", font_size=20, color=ACCENT_RED, bold=True)
add_code_block(slide, 1.0, 1.8, 5.0, 2.2,
    "# BUGGY: zeros out BOTH frames AND actions\n"
    "cfg_mask = rand(B) < 0.1\n"
    "z_cond = where(mask, zeros, z_cond)  # zero frames\n"
    "actions = where(mask, NULL_TOKEN, actions)  # zero actions!\n"
    "\n"
    "# At inference with CFG:\n"
    "v_uncond = model(z, t, z_cond=zeros, actions=NULL)\n"
    "v_cond   = model(z, t, z_cond=real,  actions=real)\n"
    "v = v_uncond + 1.5 * (v_cond - v_uncond)\n"
    "# ^ v_uncond has NO action info, so CFG amplifies\n"
    "#   both frame AND action conditioning together",
    font_size=10)

add_box(slide, 7.0, 1.2, 5.5, 3.2, fill_color=DARK_GREEN_BG, border_color=ACCENT_GREEN)
add_textbox(slide, 7.2, 1.3, 5.0, 0.4, "The Fix (GameNGen design)", font_size=20, color=ACCENT_GREEN, bold=True)
add_code_block(slide, 7.2, 1.8, 5.0, 2.2,
    "# FIXED: only drop frames, NEVER actions\n"
    "cfg_mask = rand(B) < 0.1\n"
    "z_cond = where(mask, zeros, z_cond)  # zero frames\n"
    "# actions unchanged! always pass real actions\n"
    "\n"
    "# At inference with CFG:\n"
    "v_uncond = model(z, t, z_cond=zeros, actions=REAL)\n"
    "v_cond   = model(z, t, z_cond=real,  actions=REAL)\n"
    "v = v_uncond + 1.5 * (v_cond - v_uncond)\n"
    "# ^ CFG only amplifies frame conditioning\n"
    "#   action signal preserved in both paths",
    font_size=10)

add_box(slide, 0.8, 4.7, 12.0, 2.5, border_color=ACCENT_ORANGE)
add_textbox(slide, 1.0, 4.8, 11.0, 0.5, "Observed Symptoms", font_size=20, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, 1.0, 5.3, 5.5, 1.5, [
    "Bird randomly disappeared when approaching pipes",
    "v_uncond (no context, no action) = generic 'average frame'",
    "Bird could be anywhere in the average frame",
    "CFG amplification destructively interfered with bird position",
], font_size=14, spacing=Pt(4))
add_bullet_list(slide, 7.0, 5.3, 5.5, 1.5, [
    "After fix: bird stays visible and stable",
    "Key insight from GameNGen paper (Section 3.2):",
    "\"Unconditional\" means \"no visual memory\" not \"no info\"",
    "The model should ALWAYS know what action was taken",
], font_size=14, spacing=Pt(4))


# ============================================================
# SLIDE: Bug #3 — Episode Boundary Wrap-Around
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.3, 8, 0.6, "Bug: Episode Boundary Wrap-Around", font_size=32, color=ACCENT_RED, bold=True)

add_box(slide, 0.8, 1.2, 5.5, 2.5, fill_color=DARK_RED_BG, border_color=ACCENT_RED)
add_textbox(slide, 1.0, 1.3, 5.0, 0.4, "The Bug", font_size=20, color=ACCENT_RED, bold=True)
add_bullet_list(slide, 1.0, 1.8, 5.0, 1.5, [
    "Episode N, step 0: context wraps to Episode N-1's LAST k frames",
    "Context: dead bird (crashed state) from previous episode",
    "Target: fresh bird at starting position (new episode)",
    "Model learns: dead_bird + action -> fresh_bird (impossible!)",
    "~3-4K contaminated samples (~2.4% of dataset)",
], font_size=13, spacing=Pt(4))

add_box(slide, 7.0, 1.2, 5.5, 2.5, fill_color=DARK_GREEN_BG, border_color=ACCENT_GREEN)
add_textbox(slide, 7.2, 1.3, 5.0, 0.4, "The Fix", font_size=20, color=ACCENT_GREEN, bold=True)
add_code_block(slide, 7.2, 1.8, 5.0, 1.3,
    "# OLD: only first episode starts at step k\n"
    "start_step = k if ep_idx == 0 else 0\n"
    "\n"
    "# FIXED: ALL episodes start at step k\n"
    "start_step = k  # require k prior frames from SAME episode",
    font_size=12)
add_textbox(slide, 7.2, 3.3, 5.0, 0.3, "Wastes k frames/episode but ensures physical causality", font_size=12, color=MED_GRAY)

# Other data pipeline bugs
add_box(slide, 0.8, 4.0, 12.0, 3.2, border_color=ACCENT_BLUE)
add_textbox(slide, 1.0, 4.1, 11.0, 0.5, "Other Data Pipeline Issues Found", font_size=20, color=ACCENT_BLUE, bold=True)

issues = [
    ("Inference normalization", "VAE outputs [-1,1] but code clamped to [0,1] — all negatives clipped", "Fixed: (x+1)/2"),
    ("Initial reset frame", "Frame 0 (before any action) was never captured", "Fixed: save after reset"),
    ("Train/test step mismatch", "Training: 50 steps. Inference: 20 steps. Different discretization", "Fixed: use 50 everywhere"),
    ("Latent bounds", "Generated latents could drift 10+ std from mean, garbage output", "Fixed: clamp to [-4, 4]"),
]
for i, (name, desc, fix) in enumerate(issues):
    y = 4.7 + i * 0.6
    add_textbox(slide, 1.2, y, 2.5, 0.4, name, font_size=12, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, 3.8, y, 5.5, 0.4, desc, font_size=12, color=LIGHT_GRAY)
    add_textbox(slide, 9.5, y, 3.0, 0.4, fix, font_size=12, color=ACCENT_GREEN)


# ============================================================
# SLIDE: Design Decision — Action Class Imbalance
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.6, "Design Decision: Action Class Imbalance", font_size=32, color=ACCENT_ORANGE, bold=True)

add_box(slide, 0.8, 1.2, 5.5, 3.0, border_color=ACCENT_RED)
add_textbox(slide, 1.0, 1.3, 5.0, 0.4, "The Problem", font_size=20, color=ACCENT_RED, bold=True)
add_bullet_list(slide, 1.0, 1.8, 5.0, 1.0, [
    "Action 0 (no-flap): 94.6% of training data",
    "Action 1 (flap): 5.4% — only 8,931 samples",
    "Model achieves low loss by IGNORING the action entirely",
], font_size=14, spacing=Pt(4))

add_textbox(slide, 1.0, 3.0, 5.0, 0.4, "Diagnostic: action_diff", font_size=16, color=ACCENT_ORANGE, bold=True)
add_code_block(slide, 1.0, 3.4, 5.0, 0.7,
    "# Same input, opposite actions — how different?\n"
    "action_diff = |model(ctx, action=0) - model(ctx, action=1)|\n"
    "# Result: ~0.01-0.03 (frame_diff was ~0.17)\n"
    "# Model produces IDENTICAL outputs regardless of action!",
    font_size=10)

add_box(slide, 7.0, 1.2, 5.5, 3.0, border_color=ACCENT_GREEN)
add_textbox(slide, 7.2, 1.3, 5.0, 0.4, "The Fix: Loss Reweighting", font_size=20, color=ACCENT_GREEN, bold=True)
add_code_block(slide, 7.2, 1.8, 5.0, 1.2,
    "# Per-sample weight based on target action\n"
    "weights = where(action == 1, 17.7, 1.0)\n"
    "flow_loss = sum(mse * weights) / sum(weights)\n"
    "\n"
    "# Weight = inverse frequency ratio\n"
    "# Normalized by weight sum (not batch size)\n"
    "# Stable loss magnitude across batches",
    font_size=11)

add_bullet_list(slide, 7.2, 3.2, 5.0, 0.8, [
    "Weights computed automatically from dataset during encoding",
    "Same approach for done_pos_weight (~28x for terminal frames)",
], font_size=13, spacing=Pt(4))

add_box(slide, 0.8, 4.5, 12.0, 2.8, border_color=ACCENT_PURPLE)
add_textbox(slide, 1.0, 4.6, 11.0, 0.5, "The Lazy Extrapolation Problem", font_size=20, color=ACCENT_PURPLE, bold=True)
add_bullet_list(slide, 1.0, 5.1, 11.0, 2.0, [
    "Even after reweighting, the model had a shortcut: with k=16 context frames, it could extrapolate the bird's trajectory from visual momentum alone",
    "Dense temporal context explains 95%+ of variance. The action explains the remaining few percent — easy to ignore",
    "Fundamental tension: MORE context helps scene understanding but WEAKENS incentive to use action signal",
    "This is why we eventually reduced context from k=16 to k=8 and then to k=4 — trading temporal context for action sensitivity",
], font_size=13, spacing=Pt(4))


# ============================================================
# SLIDE: Design Decision — Death Frames Semantics
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.6, "Design Decision: Death Frame Semantics", font_size=32, color=ACCENT_ORANGE, bold=True)

add_box(slide, 0.8, 1.2, 5.5, 3.5, fill_color=DARK_RED_BG, border_color=ACCENT_RED)
add_textbox(slide, 1.0, 1.3, 5.0, 0.4, "Original: 5 Death Frames (all terminated=True)", font_size=16, color=ACCENT_RED, bold=True)
add_bullet_list(slide, 1.0, 1.8, 5.0, 2.5, [
    "Collision frame: terminated=True, NO overlay (looks normal!)",
    "5 overlay frames: terminated=True, game-over sprite",
    "Problems:",
    "  Done head learns shortcut: detect overlay pixels = done",
    "  6 terminated frames per episode inflates done count",
    "  Pollutes context window (k=8: 62% death frames)",
    "  Wasted generation capacity — model generates overlays",
    "    but inference stops at first done=True anyway",
], font_size=12, spacing=Pt(3))

add_box(slide, 7.0, 1.2, 5.5, 3.5, fill_color=DARK_GREEN_BG, border_color=ACCENT_GREEN)
add_textbox(slide, 7.2, 1.3, 5.0, 0.4, "Fixed: 1 Crash Frame + 1 Overlay Frame", font_size=16, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, 7.2, 1.8, 5.0, 2.5, [
    "Crash frame: terminated=True (the real collision)",
    "1 overlay frame: terminated=FALSE, post_terminal=True",
    "Clean semantic split:",
    "  Done head learns to detect COLLISION, not overlay",
    "  Model learns crash -> overlay as normal prediction",
    "  At inference: done fires, 1 more frame generated",
    "  Model produces its OWN game-over visual, then freeze",
    "  All generated pixels from the world model itself",
], font_size=12, spacing=Pt(3))

add_box(slide, 0.8, 5.0, 12.0, 2.2, border_color=ACCENT_BLUE)
add_textbox(slide, 1.0, 5.1, 11.0, 0.5, "Done Head: Training/Inference Mismatch", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, 1.0, 5.6, 5.5, 1.5, [
    "Training: done head sees z_t at random t ~ U(0,1)",
    "Inference: done head only queried at t=1 (clean sample)",
    "Predictions at low t (noisy intermediates) never used",
], font_size=13, spacing=Pt(3))
add_bullet_list(slide, 7.0, 5.6, 5.5, 1.5, [
    "Solution: weight done loss by w(t) = 5 * t^4",
    "At t=0.1: weight = 0.0005 (negligible)",
    "At t=0.9: weight = 3.28 (dominant)",
    "Normalized to mean-1 under U(0,1) for stable scaling",
], font_size=13, spacing=Pt(3))


# ============================================================
# SLIDE: Design Decision — Action Conditioning Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.6, "Design Decision: How Should Actions Enter the Model?", font_size=32, color=ACCENT_ORANGE, bold=True)

# Three approaches compared
add_box(slide, 0.5, 1.2, 3.8, 3.2, fill_color=DARK_RED_BG, border_color=ACCENT_RED)
add_textbox(slide, 0.7, 1.3, 3.4, 0.4, "V1: AdaGN Only", font_size=16, color=ACCENT_RED, bold=True)
add_bullet_list(slide, 0.7, 1.8, 3.4, 2.0, [
    "All k+1 actions flattened into one vector",
    "Global scale/shift via AdaGN at every layer",
    "Spatially uniform — same modulation everywhere",
    "Model must learn to route global signal",
    "to bird region through conv filters",
    "Target action buried in flattened vector",
    "Result: model mostly ignores actions",
], font_size=11, spacing=Pt(3))

add_box(slide, 4.7, 1.2, 3.8, 3.2, border_color=ACCENT_BLUE)
add_textbox(slide, 4.9, 1.3, 3.4, 0.4, "V2: Spatial Broadcast Only", font_size=16, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, 4.9, 1.8, 3.4, 2.0, [
    "Target action embedded + broadcast to (B,16,H,W)",
    "Concatenated as input channels",
    "First conv sees action spatially",
    "Context actions DISCARDED entirely",
    "Only first layer processes the action",
    "No deep action injection",
    "Result: better, but gravity still broken",
], font_size=11, spacing=Pt(3))

add_box(slide, 9.0, 1.2, 3.8, 3.2, fill_color=DARK_GREEN_BG, border_color=ACCENT_GREEN)
add_textbox(slide, 9.2, 1.3, 3.4, 0.4, "V3: Dual Path (current)", font_size=16, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, 9.2, 1.8, 3.4, 2.0, [
    "BOTH spatial + AdaGN simultaneously",
    "Target action: broadcast spatially at input",
    "All k+1 actions: AdaGN at every layer",
    "Action impossible to ignore — woven into",
    "input tensor AND deep layer modulations",
    "Shared Embedding(2, 128) for both paths",
    "Result: improved, but gravity still a problem",
], font_size=11, spacing=Pt(3))

add_box(slide, 0.8, 4.7, 12.0, 2.5, border_color=ACCENT_ORANGE)
add_textbox(slide, 1.0, 4.8, 11.0, 0.5, "Why Is This So Hard?", font_size=20, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, 1.0, 5.3, 11.0, 1.5, [
    "Flappy Bird has a BINARY action space (flap / no-flap). Cross-attention (GameNGen's approach) is overkill.",
    "The challenge isn't representing the action — it's making the model USE it over the easier alternative of extrapolating from temporal context",
    "AdaGN is inherently global: same scale/shift everywhere. But 'flap' means 'apply upward velocity to the bird specifically' — a spatially localized effect",
    "Zero-initialized projection means action conditioning literally starts as a no-op and must be learned from scratch",
], font_size=13, spacing=Pt(4))


# ============================================================
# SLIDE: Failed Experiment — Residual Prediction
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.3, 8, 0.6, "Failed Experiment: Residual Prediction", font_size=32, color=ACCENT_RED, bold=True)
add_textbox(slide, 9.5, 0.4, 3, 0.4, "REVERTED", font_size=20, color=ACCENT_RED, bold=True)

add_box(slide, 0.8, 1.2, 5.5, 2.8, border_color=ACCENT_BLUE)
add_textbox(slide, 1.0, 1.3, 5.0, 0.4, "The Idea", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, 1.0, 1.8, 5.0, 2.0, [
    "Instead of predicting z_target from N(0,1)...",
    "Predict the DELTA: z_target - z_last (frame difference)",
    "Hypothesis: frame deltas capture dynamics (velocity, etc.)",
    "Should make the model focus on what CHANGED",
    "Similar to residual connections — predict the update",
], font_size=13, spacing=Pt(4))

add_box(slide, 7.0, 1.2, 5.5, 2.8, fill_color=DARK_RED_BG, border_color=ACCENT_RED)
add_textbox(slide, 7.2, 1.3, 5.0, 0.4, "Why It Failed Catastrophically", font_size=20, color=ACCENT_RED, bold=True)
add_bullet_list(slide, 7.2, 1.8, 5.0, 2.0, [
    "Frame deltas are NEAR ZERO (consecutive frames very similar)",
    "z_1 = z_target - z_last ~ 0",
    "But z_0 ~ N(0,1) has magnitude ~ 1.0",
    "So v_target = z_1 - z_0 ~ -z_0 (just negative noise!)",
    "The actual frame-delta signal is LOST in the noise",
], font_size=13, spacing=Pt(4))

add_box(slide, 0.8, 4.3, 12.0, 3.0, border_color=ACCENT_ORANGE)
add_textbox(slide, 1.0, 4.4, 11.0, 0.5, "The Result", font_size=20, color=ACCENT_ORANGE, bold=True)
add_code_block(slide, 1.0, 4.9, 5.5, 1.5,
    "# What we expected:\n"
    "# Model learns meaningful velocity field\n"
    "# for small frame-to-frame changes\n"
    "\n"
    "# What actually happened:\n"
    "# v_target ≈ 0 - z_0 = -z_0\n"
    "# Model learns to predict negative of input noise\n"
    "# Output: totally unstable garbage frames",
    font_size=11)
add_bullet_list(slide, 7.0, 5.0, 5.5, 1.5, [
    "Lesson: flow matching from N(0,1) requires the TARGET to have large magnitude",
    "Near-zero targets make the flow degenerate — the model just learns to cancel the noise",
    "Would need a different source distribution (not N(0,1)) or a different parameterization",
    "Reverted immediately after observing unstable inference",
], font_size=13, spacing=Pt(4))


# ============================================================
# SLIDE: Failed Experiment — Filter Action=0
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.6, "Diagnostic Experiment: Train Only on No-Flap Frames", font_size=32, color=ACCENT_ORANGE, bold=True)

add_box(slide, 0.8, 1.2, 5.5, 2.5, border_color=ACCENT_BLUE)
add_textbox(slide, 1.0, 1.3, 5.0, 0.4, "Hypothesis", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, 1.0, 1.8, 5.0, 1.5, [
    "If we train ONLY on action=0 (no-flap) frames...",
    "The bird should ALWAYS fall due to gravity",
    "No conflicting flap signals to confuse the model",
    "If it works: action conditioning is the bottleneck",
    "If it doesn't work: something more fundamental is wrong",
], font_size=14, spacing=Pt(4))

add_box(slide, 7.0, 1.2, 5.5, 2.5, fill_color=DARK_RED_BG, border_color=ACCENT_RED)
add_textbox(slide, 7.2, 1.3, 5.0, 0.4, "Result: Bird STILL Doesn't Fall", font_size=20, color=ACCENT_RED, bold=True)
add_bullet_list(slide, 7.2, 1.8, 5.0, 1.5, [
    "Even with 100% no-flap training data, no gravity",
    "Bird hovers near pipe center, never falls",
    "Rules out action conditioning as sole cause",
    "Points to a fundamental optimization problem",
], font_size=14, spacing=Pt(4))

add_box(slide, 0.8, 4.0, 12.0, 3.2, border_color=ACCENT_PURPLE)
add_textbox(slide, 1.0, 4.1, 11.0, 0.5, "Root Cause Analysis: The Bird Is Too Small", font_size=20, color=ACCENT_PURPLE, bold=True)
add_bullet_list(slide, 1.0, 4.7, 5.5, 2.0, [
    "The bird occupies ~1-2% of the latent spatial area",
    "MSE loss over full frame: 98% of gradient from background",
    "Background is EASY (pipes scroll predictably, sky is static)",
    "Model achieves low loss by perfecting the background",
    "and predicting bird at the 'average' position",
], font_size=14, spacing=Pt(4))

add_box(slide, 7.0, 4.7, 5.5, 2.3, border_color=ACCENT_GREEN)
add_textbox(slide, 7.2, 4.8, 5.0, 0.4, "Information Theory Check", font_size=18, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, 7.2, 5.3, 5.0, 1.5, [
    "Model: ~3.4M params x 2 bits/param = 6.8M bits",
    "Flappy Bird state: ~100K bits (generous estimate)",
    "Model has ~70x more capacity than needed",
    "NOT a capacity bottleneck — it's OPTIMIZATION",
    "MSE doesn't provide enough gradient for bird region",
], font_size=14, spacing=Pt(4))


# ============================================================
# SLIDE: Solution — Auxiliary Dynamics Head
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.6, "Solution: Auxiliary Dynamics Head", font_size=32, color=ACCENT_GREEN, bold=True)

add_box(slide, 0.8, 1.2, 5.5, 3.0, border_color=ACCENT_BLUE)
add_textbox(slide, 1.0, 1.3, 5.0, 0.4, "The Idea", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, 1.0, 1.8, 5.0, 2.0, [
    "Add MLP head at bottleneck that predicts game state:",
    "  [dy, vel_y, dx] = bird-to-gap dist, velocity, pipe dist",
    "These are available in training data (processed_obs)",
    "MSE loss on this prediction, weight = 0.1",
    "Crucially: NOT detached from backbone",
    "Gradients flow back into the entire encoder",
], font_size=14, spacing=Pt(4))

add_box(slide, 7.0, 1.2, 5.5, 3.0, border_color=ACCENT_GREEN)
add_textbox(slide, 7.2, 1.3, 5.0, 0.4, "Why This Should Work", font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, 7.2, 1.8, 5.0, 2.0, [
    "Forces bottleneck to learn bird position/velocity features",
    "Direct gradient signal for bird dynamics (not buried in MSE)",
    "Contrast with done head: done IS detached",
    "  (done is passive reader, doesn't distort flow features)",
    "Dynamics head IS NOT detached:",
    "  (want to actively shape representations to know bird state)",
], font_size=14, spacing=Pt(4))

add_box(slide, 0.8, 4.5, 12.0, 2.7, border_color=ACCENT_ORANGE)
add_textbox(slide, 1.0, 4.6, 11.0, 0.5, "Architecture", font_size=20, color=ACCENT_ORANGE, bold=True)
add_code_block(slide, 1.0, 5.1, 5.5, 1.5,
    "# Dynamics head (NOT detached from backbone)\n"
    "dynamics_head = Sequential(\n"
    "    AdaptiveAvgPool2d(1),  # (B, 256, 9, 16) -> (B, 256)\n"
    "    Flatten(),\n"
    "    Linear(256, 64), SiLU(),\n"
    "    Linear(64, 3),         # -> [dy, vel_y, dx]\n"
    ")",
    font_size=11)
add_code_block(slide, 7.0, 5.1, 5.5, 1.5,
    "# In loss function:\n"
    "dynamics_loss = MSE(dynamics_pred, game_states)\n"
    "\n"
    "total = flow_loss\n"
    "     + 0.1 * done_loss    # detached\n"
    "     + 0.1 * dynamics_loss # NOT detached\n"
    "\n"
    "# Current experiment — awaiting results",
    font_size=11)


# ============================================================
# SLIDE: Design Decision — Noise Augmentation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.6, "Design Decision: Noise Augmentation for Autoregressive Stability", font_size=30, color=ACCENT_ORANGE, bold=True)

add_box(slide, 0.8, 1.2, 5.5, 2.5, border_color=ACCENT_RED)
add_textbox(slide, 1.0, 1.3, 5.0, 0.4, "The Problem", font_size=20, color=ACCENT_RED, bold=True)
add_bullet_list(slide, 1.0, 1.8, 5.0, 1.5, [
    "At inference, model conditions on its OWN outputs",
    "Training sees only CLEAN (ground truth) conditioning",
    "Small errors compound: drift, artifacts, collapse",
    "Quality degrades within seconds of generation",
], font_size=14, spacing=Pt(4))

add_box(slide, 7.0, 1.2, 5.5, 2.5, border_color=ACCENT_GREEN)
add_textbox(slide, 7.2, 1.3, 5.0, 0.4, "The Solution (from GameNGen)", font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, 7.2, 1.8, 5.0, 1.5, [
    "Add random Gaussian noise to conditioning during training",
    "Model learns to handle imperfect, noisy context",
    "Naturally robust to own prediction errors at inference",
    "Tell model the noise level via discrete embedding",
], font_size=14, spacing=Pt(4))

add_box(slide, 0.8, 4.0, 12.0, 3.2, border_color=ACCENT_BLUE)
add_textbox(slide, 1.0, 4.1, 11.0, 0.5, "Design Choices and Tradeoffs", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, 1.0, 4.7, 5.5, 2.0, [
    "Discretized bins (10) vs continuous: gives model clear signal",
    "max_aug_std = 0.5: tried reducing to 0.25 (no improvement)",
    "GameNGen uses 0.7 — higher noise = more robustness",
    "At inference: set aug_level=0 ('this conditioning is clean')",
], font_size=13, spacing=Pt(4))
add_bullet_list(slide, 7.0, 4.7, 5.5, 2.0, [
    "aug_level enters model through Embedding(10, 128) -> SiLU",
    "Concatenated with time and action in AdaGN",
    "Also tried aug_level=8 at inference (match training median)",
    "Settled on aug_level=0 — cleaner outputs for first frame",
    "The augmentation also helps with the done head:",
    "  noisy context still lets model predict termination",
], font_size=13, spacing=Pt(4))


# ============================================================
# SLIDE: Design Decision — Behavioral Diversity
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.6, "Design Decision: Data Collection & Behavioral Diversity", font_size=30, color=ACCENT_ORANGE, bold=True)

add_box(slide, 0.8, 1.2, 5.5, 2.5, border_color=ACCENT_RED)
add_textbox(slide, 1.0, 1.3, 5.0, 0.4, "Problem: Expert Data Is Boring", font_size=20, color=ACCENT_RED, bold=True)
add_bullet_list(slide, 1.0, 1.8, 5.0, 1.5, [
    "Optimal PPO agent plays perfectly — narrow state coverage",
    "Bird always near pipe gaps, never extreme heights",
    "Almost never crashes in interesting ways",
    "Model has no training signal for 'what if bird is too high?'",
], font_size=14, spacing=Pt(4))

add_box(slide, 7.0, 1.2, 5.5, 2.5, border_color=ACCENT_GREEN)
add_textbox(slide, 7.2, 1.3, 5.0, 0.4, "Solution: Hijacked PPO Agent", font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, 7.2, 1.8, 5.0, 1.5, [
    "p_stim: probability of forcing a flap (when agent says don't)",
    "p_freeze: probability of blocking a flap (when agent says do)",
    "Grid of 17 configurations: (0, 0) to (0.05, 0.5)",
    "Creates birds flying too high, crashing into ground, etc.",
], font_size=14, spacing=Pt(4))

add_box(slide, 0.8, 4.0, 12.0, 3.2, border_color=ACCENT_PURPLE)
add_textbox(slide, 1.0, 4.1, 11.0, 0.5, "Data Evolution Across Runs", font_size=20, color=ACCENT_PURPLE, bold=True)

runs = [
    ("1/29 Run", "16 configs, 50 runs each", "141K frames, 1.3hr", "Bird never extreme heights"),
    ("1/31 Run", "17 configs, up to 330 runs", "467K frames, 4.3hr", "3,135 episodes, much more diverse"),
    ("Future", "More p_stim/p_freeze spread", "Target 1M+ frames", "Cover edge cases: ground crash, ceiling, etc."),
]
for i, (run, config, data, note) in enumerate(runs):
    y = 4.7 + i * 0.7
    add_textbox(slide, 1.2, y, 1.5, 0.4, run, font_size=13, color=ACCENT_PURPLE, bold=True)
    add_textbox(slide, 2.8, y, 3.0, 0.4, config, font_size=13, color=LIGHT_GRAY)
    add_textbox(slide, 5.8, y, 2.5, 0.4, data, font_size=13, color=LIGHT_GRAY)
    add_textbox(slide, 8.5, y, 4.0, 0.4, note, font_size=13, color=MED_GRAY)


# ============================================================
# SLIDE: Design Decision — Bird-Weighted VAE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.6, "Design Decision: Bird-Weighted VAE Loss", font_size=32, color=ACCENT_ORANGE, bold=True)

add_box(slide, 0.8, 1.2, 5.5, 2.8, border_color=ACCENT_RED)
add_textbox(slide, 1.0, 1.3, 5.0, 0.4, "The Problem", font_size=20, color=ACCENT_RED, bold=True)
add_bullet_list(slide, 1.0, 1.8, 5.0, 2.0, [
    "Flappy Bird frames: mostly sky, ground, pipes",
    "Bird is ~2% of the image area",
    "Standard L1 loss: 98% of gradient from background",
    "VAE happily sacrifices bird detail for background fidelity",
], font_size=14, spacing=Pt(4))

add_box(slide, 7.0, 1.2, 5.5, 2.8, border_color=ACCENT_GREEN)
add_textbox(slide, 7.2, 1.3, 5.0, 0.4, "The Solution: 10x Bird Weight", font_size=20, color=ACCENT_GREEN, bold=True)
add_code_block(slide, 7.2, 1.8, 5.0, 2.0,
    "# Detect bird by color (it's orange)\n"
    "bird_mask = (r>0.8) & (g>0.3) & (g<0.6) & (b<0.3)\n"
    "\n"
    "# Find top, offset 23px for white head\n"
    "y_min = rows.min() - 23\n"
    "\n"
    "# Fixed x range (bird doesn't move horizontally)\n"
    "weights[..., y_min:y_max, 50:100] = 10.0\n"
    "\n"
    "# VAE loss = weighted_L1 + gradient_L1 + 0.001*KL",
    font_size=11)

add_box(slide, 0.8, 4.3, 12.0, 2.0, border_color=ACCENT_BLUE)
add_textbox(slide, 1.0, 4.4, 11.0, 0.5, "VAE Loss Composition", font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, 1.0, 4.9, 5.5, 1.0, [
    "Weighted L1: 10x on bird region for sharp reconstruction",
    "Gradient L1: preserves edges (pipes, ground line, bird outline)",
], font_size=14, spacing=Pt(3))
add_bullet_list(slide, 7.0, 4.9, 5.5, 1.0, [
    "KL divergence (w=0.001): regularizes latent space structure",
    "Also added 10x weight for game-over overlay sprite",
], font_size=14, spacing=Pt(3))

add_textbox(slide, 0.8, 6.5, 12, 0.5, "External normalization: z = (z_raw - 0.4735) / 1.5931. Decouples VAE and flow model — swap checkpoints without retuning.", font_size=13, color=MED_GRAY)


# ============================================================
# SLIDE: Engineering — Optimizations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.6, "Engineering: Optimizations & Infrastructure", font_size=32, color=ACCENT_BLUE, bold=True)

add_box(slide, 0.8, 1.2, 5.5, 3.0, border_color=ACCENT_GREEN)
add_textbox(slide, 1.0, 1.3, 5.0, 0.4, "What Worked", font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, 1.0, 1.8, 5.0, 2.0, [
    "torch.compile(mode='reduce-overhead') — CUDA graphs",
    "bfloat16 autocast — ~2x memory, ~30% speedup",
    "Gradient checkpointing — ~50% memory at ~25% compute cost",
    "DDP with NCCL — multi-GPU scaling",
    "Persistent workers + prefetch_factor=2",
    "Pre-encode all latents (was encoding on-the-fly earlier!)",
    "torch.inference_mode() at test time",
], font_size=13, spacing=Pt(3))

add_box(slide, 7.0, 1.2, 5.5, 3.0, fill_color=DARK_RED_BG, border_color=ACCENT_RED)
add_textbox(slide, 7.2, 1.3, 5.0, 0.4, "What Didn't Work / Reverted", font_size=20, color=ACCENT_RED, bold=True)
add_bullet_list(slide, 7.2, 1.8, 5.0, 2.0, [
    "channels_last memory format — crashed on rank-5 tensors",
    "  (past_frames is [B, k, C, H, W] = rank 5, needs rank 4)",
    "ConvTranspose2d for upsampling — checkerboard artifacts",
    "  -> Replaced with interpolate + conv",
    "Running encode_vod.py every training epoch (very slow!)",
    "  -> Pre-encode once, load tensors during training",
], font_size=13, spacing=Pt(3))

add_box(slide, 0.8, 4.5, 12.0, 2.8, border_color=ACCENT_PURPLE)
add_textbox(slide, 1.0, 4.6, 11.0, 0.5, "Automation: Fewer Manual Steps", font_size=20, color=ACCENT_PURPLE, bold=True)
add_bullet_list(slide, 1.0, 5.1, 5.5, 1.5, [
    "Auto latent statistics: mean/std computed after VAE training",
    "  -> Saved to config.json, loaded by encode_vod.py",
    "Auto class weights: action_weight and done_pos_weight",
    "  -> Computed during encoding, saved to encode_config.json",
    "  -> Loaded by train_ngen.py at startup",
], font_size=13, spacing=Pt(3))
add_bullet_list(slide, 7.0, 5.1, 5.5, 1.5, [
    "Auto resume: training resumes from latest.pt checkpoint",
    "  -> Optimizer + scheduler + epoch state restored",
    "Config saved once: model_config + train_config + paths",
    "  -> Reconstruct exact training setup from run directory",
    "Run directory broadcast via DDP for multi-GPU consistency",
], font_size=13, spacing=Pt(3))


# ============================================================
# SLIDE: Full Timeline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.6, "Full Experiment Timeline", font_size=32, color=WHITE, bold=True)

timeline = [
    ("1/16", "First ngen run", "Encoding on-the-fly, ~40 epochs, no latent-vod", ACCENT_BLUE, "Baseline"),
    ("1/20", "Added done head", "BCE loss on termination prediction", ACCENT_BLUE, "Feature"),
    ("1/23", "Added action weight", "17x loss reweight for flap samples", ACCENT_GREEN, "Fix"),
    ("1/26", "Added latent-vod", "Pre-encode everything, huge speedup", ACCENT_GREEN, "Perf"),
    ("1/27", "Optimization sweep", "torch.compile, bf16, grad checkpoint, grad clip, LR sched", ACCENT_GREEN, "Perf"),
    ("1/29", "New data: 141K frames", "16 hijacking configs, VAE retrained", ACCENT_BLUE, "Data"),
    ("1/30", "Full training run", "k=8, CFG, attention at bottleneck, 2 H200s", ACCENT_PURPLE, "Run"),
    ("", "", "Result: HORRIBLE — worse than before, ignores gravity", ACCENT_RED, ""),
    ("1/30", "Found CFG bug", "Was dropping actions in unconditional path", ACCENT_RED, "Bug"),
    ("1/31", "Death frames fix", "5 -> 1, done head t^4 reweighting", ACCENT_ORANGE, "Fix"),
    ("1/31", "New data: 467K frames", "3,135 episodes, 17 configs, 4.3hr", ACCENT_BLUE, "Data"),
    ("2/01", "New VAE + encoding", "hidden_ch=32, game-over sprite support", ACCENT_BLUE, "Train"),
    ("2/02", "Ngen run 2", "Bird disappears near pipes, no gravity", ACCENT_RED, "Fail"),
    ("2/06", "Spatial action cond", "Broadcast target action as input channels", ACCENT_GREEN, "Arch"),
    ("2/06", "Dual-path (spatial+AdaGN)", "Actions slightly better, gravity still fails", ACCENT_ORANGE, "Run"),
    ("2/26", "Residual prediction", "Predict z_target - z_last. REVERTED (degenerate)", ACCENT_RED, "Fail"),
    ("2/26", "Filter action=0 test", "Train only no-flap. Bird still doesn't fall!", ACCENT_RED, "Diag"),
    ("2/26", "Info theory analysis", "6.8M bits capacity vs 100K needed — optimization problem", ACCENT_PURPLE, "Insight"),
    ("2/26", "Dynamics head", "Predict [dy, vel_y, dx] from bottleneck (not detached)", ACCENT_GREEN, "Current"),
]

for i, (date, name, desc, color, tag) in enumerate(timeline):
    y = 1.0 + i * 0.325
    if date:
        add_textbox(slide, 0.5, y, 0.7, 0.3, date, font_size=9, color=MED_GRAY, font_name="Courier New")
    if tag:
        add_textbox(slide, 1.2, y, 0.7, 0.3, tag, font_size=9, color=color, bold=True)
    add_textbox(slide, 2.0, y, 2.5, 0.3, name, font_size=10, color=color, bold=bool(date))
    add_textbox(slide, 4.5, y, 8.5, 0.3, desc, font_size=10, color=LIGHT_GRAY if date else ACCENT_RED)


# ============================================================
# SLIDE: What's Working, What's Not
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.6, "Current Status: What Works, What Doesn't", font_size=32, color=WHITE, bold=True)

add_box(slide, 0.8, 1.2, 5.5, 3.0, border_color=ACCENT_GREEN)
add_textbox(slide, 1.0, 1.3, 5.0, 0.4, "Working", font_size=22, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, 1.0, 1.8, 5.0, 2.0, [
    "Visually sharp frame generation (VAE is excellent)",
    "Pipes scroll correctly, background moves",
    "Overall scene structure is coherent",
    "Done head detects termination near crash frames",
    "Actions produce visible (if subtle) differences",
    "Interactive play loop runs at ~10 FPS",
    "Training infrastructure is solid and automated",
], font_size=13, spacing=Pt(4))

add_box(slide, 7.0, 1.2, 5.5, 3.0, fill_color=DARK_RED_BG, border_color=ACCENT_RED)
add_textbox(slide, 7.2, 1.3, 5.0, 0.4, "Not Working", font_size=22, color=ACCENT_RED, bold=True)
add_bullet_list(slide, 7.2, 1.8, 5.0, 2.0, [
    "Bird doesn't respect gravity (floats/hovers)",
    "Bird occasionally disappears near pipes",
    "Action response is visible but not crisp",
    "Model predicts 'average' bird position instead of physics",
    "Long-horizon generation drifts after 30+ seconds",
], font_size=13, spacing=Pt(4))

add_box(slide, 0.8, 4.5, 12.0, 2.8, border_color=ACCENT_PURPLE)
add_textbox(slide, 1.0, 4.6, 11.0, 0.5, "Key Insight: The Bottleneck Is Optimization, Not Capacity", font_size=20, color=ACCENT_PURPLE, bold=True)
add_bullet_list(slide, 1.0, 5.1, 11.0, 2.0, [
    "Model has 70x more capacity than needed (6.8M bits vs ~100K bits for full game state)",
    "The bird occupies ~1-2% of latent spatial area — MSE gradient is dominated by background",
    "Background is easy to predict (pipes scroll linearly, sky is static). Bird dynamics are hard (acceleration, collision, action-dependent).",
    "The model takes the path of least resistance: perfect background, average bird position = low MSE without learning physics",
    "This is why the dynamics head is promising: it provides DIRECT gradient signal for bird position/velocity, bypassing the MSE bottleneck",
], font_size=13, spacing=Pt(4))


# ============================================================
# SLIDE: GameNGen Comparison
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.6, "GameNGen vs This Project", font_size=32, color=WHITE, bold=True)

rows = [
    ("Base model", "Stable Diffusion 1.4 (~860M params)", "Custom ResUNet (~3.4M params)"),
    ("Training approach", "Fine-tune pretrained", "Train from scratch"),
    ("Generative method", "DDIM denoising diffusion", "Flow matching + Euler integration"),
    ("Data scale", "900M frames", "467K frames (2000x less)"),
    ("Context frames", "64 (3.2s)", "4 (0.13s)"),
    ("Action conditioning", "Cross-attention", "Dual-path: spatial + AdaGN"),
    ("CFG", "Drop frames only", "Drop frames only (after bug fix)"),
    ("Noise augmentation", "0.7 max, 10 bins", "0.5 max, 10 bins"),
    ("Inference steps", "4 DDIM", "50 Euler (planned: reflow to 8)"),
    ("Game", "DOOM", "Flappy Bird"),
    ("Done detection", "N/A (DOOM doesn't end)", "Auxiliary head + t^4 weighting"),
    ("Dynamics head", "N/A", "Predict [dy, vel_y, dx] from bottleneck"),
]

add_textbox(slide, 1.0, 1.0, 3.5, 0.4, "Aspect", font_size=14, color=ACCENT_BLUE, bold=True)
add_textbox(slide, 4.5, 1.0, 4.0, 0.4, "GameNGen", font_size=14, color=ACCENT_BLUE, bold=True)
add_textbox(slide, 8.8, 1.0, 4.0, 0.4, "This Project", font_size=14, color=ACCENT_BLUE, bold=True)

for i, (aspect, gamengen, ours) in enumerate(rows):
    y = 1.5 + i * 0.46
    bg = BG_MED if i % 2 == 0 else BG_DARK
    add_box(slide, 0.8, y, 11.8, 0.42, fill_color=bg)
    add_textbox(slide, 1.0, y + 0.02, 3.3, 0.35, aspect, font_size=12, color=LIGHT_GRAY, bold=True)
    add_textbox(slide, 4.5, y + 0.02, 4.0, 0.35, gamengen, font_size=12, color=LIGHT_GRAY)
    add_textbox(slide, 8.8, y + 0.02, 4.0, 0.35, ours, font_size=12, color=LIGHT_GRAY)


# ============================================================
# SLIDE: Next Steps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, 0.8, 0.3, 12, 0.6, "Next Steps & Open Questions", font_size=32, color=WHITE, bold=True)

add_box(slide, 0.8, 1.2, 5.5, 3.0, border_color=ACCENT_GREEN)
add_textbox(slide, 1.0, 1.3, 5.0, 0.4, "Immediate Next Steps", font_size=20, color=ACCENT_GREEN, bold=True)
add_bullet_list(slide, 1.0, 1.8, 5.0, 2.0, [
    "Evaluate dynamics head results — does bird learn gravity?",
    "If not: increase dynamics_loss_weight, try different game state features",
    "Reflow training for faster inference (50 -> 8 steps)",
    "Decoder fine-tuning on generated latents (GameNGen approach)",
    "More data: target 1M+ frames with wider behavioral coverage",
], font_size=14, spacing=Pt(4))

add_box(slide, 7.0, 1.2, 5.5, 3.0, border_color=ACCENT_ORANGE)
add_textbox(slide, 7.2, 1.3, 5.0, 0.4, "Open Questions", font_size=20, color=ACCENT_ORANGE, bold=True)
add_bullet_list(slide, 7.2, 1.8, 5.0, 2.0, [
    "How much context is optimal? (4 vs 8 vs 16 frames)",
    "More VAE compression? 36x64 -> 16x32 to force latent structure",
    "Can the model learn collision detection from pixels alone?",
    "Is reflow enough for real-time 30 FPS, or need distillation?",
    "Would a GAN/perceptual loss help the bird region?",
], font_size=14, spacing=Pt(4))

add_box(slide, 0.8, 4.5, 12.0, 2.8, border_color=ACCENT_PURPLE)
add_textbox(slide, 1.0, 4.6, 11.0, 0.5, "The Big Takeaway", font_size=22, color=ACCENT_PURPLE, bold=True)
add_bullet_list(slide, 1.0, 5.2, 11.0, 1.5, [
    "The architecture works. The training pipeline works. The model generates recognizable Flappy Bird gameplay.",
    "The gap between 'recognizable' and 'convincing' is where all the interesting design questions live.",
    "They are questions about inductive biases, not scale: How does conditioning structure shape what the model learns to rely on?",
    "When does providing more information actually hurt by enabling shortcuts? How strong does a signal need to be before the model uses it?",
], font_size=14, spacing=Pt(5))


# ============================================================
# Save
# ============================================================
output_path = "/Users/william/Desktop/Random/flappy/flappy_world_model_presentation.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
